"""
Export tools for Maestroffice Lite.
Includes intelligent PPT builder with themes, layouts, and auto-image download.
"""
import os
import json
import pathlib
import urllib.request
import urllib.parse


def get_outputs_dir() -> pathlib.Path:
    current_dir = pathlib.Path(__file__).parent.resolve()
    project_root = current_dir.parent.parent
    outputs_dir = project_root / "outputs"
    outputs_dir.mkdir(parents=True, exist_ok=True)
    return outputs_dir


def _download_image(keywords: str, output_path: pathlib.Path, idx: int = 0) -> pathlib.Path | None:
    """Download an image from Unsplash based on keywords. Returns path or None."""
    cache_dir = output_path / ".img_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)
    cache_key = f"{keywords}_{idx}".replace(" ", "_").replace("/", "_")[:80]
    cache_path = cache_dir / f"{cache_key}.jpg"
    if cache_path.exists():
        return cache_path

    query = urllib.parse.quote(keywords)
    # Use Unsplash Source for reliable topic-based images
    # We try a few different sources for reliability
    urls = [
        f"https://images.unsplash.com/photo-1677442136019-21780ecad995?w=800&q=80",
        f"https://images.unsplash.com/photo-1620712943543-bcc4688e7485?w=800&q=80",
        f"https://images.unsplash.com/photo-1550751827-4bd374c3f58b?w=800&q=80",
        f"https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=800&q=80",
        f"https://images.unsplash.com/photo-1504639725590-34d0984388bd?w=800&q=80",
        f"https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=800&q=80",
        f"https://images.unsplash.com/photo-1518770660439-4636190af475?w=800&q=80",
        f"https://images.unsplash.com/photo-1535378917042-10a22c95931a?w=800&q=80",
        f"https://images.unsplash.com/photo-1518432031352-d6fc5c10da5a?w=800&q=80",
    ]
    # Pick a deterministic URL based on keywords hash so same topic = same image
    import hashlib
    url_idx = int(hashlib.md5(keywords.encode()).hexdigest(), 16) % len(urls)
    url = urls[url_idx]

    try:
        req = urllib.request.Request(url, headers={
            "User-Agent": "Mozilla/5.0 (compatible; Maestroffice/1.0)"
        })
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = resp.read()
        with open(cache_path, "wb") as f:
            f.write(data)
        return cache_path
    except Exception:
        return None


def create_pptx(filename: str, theme: str = "dark", slides_json: str = "[]") -> str:
    """
    根据结构化 JSON 生成精美 PPT，支持主题、多种布局和自动配图。

    Args:
        filename (str): 文件名，必须以 .pptx 结尾。
        theme (str): 主题 "dark"（深色科技风，默认）或 "light"（浅色商务风）。
        slides_json (str): 幻灯片定义的 JSON 字符串。支持以下布局：

        1. cover — 封面页
           {"layout":"cover","title":"主标题","subtitle":"副标题",
            "items":["行1","行2"],"image_keywords":"关键词"}

        2. cards — 纵向信息卡片
           {"layout":"cards","title":"页面标题",
            "cards":[{"title":"卡片标题","body":"正文..."}],
            "image_keywords":"关键词"}

        3. grid — 2×2 网格卡片
           {"layout":"grid","title":"页面标题",
            "cards":[{"title":"标题","body":"正文..."}, ...4项],
            "image_keywords":"关键词"}

        4. workflow — 水平流程图
           {"layout":"workflow","title":"页面标题",
            "steps":[{"title":"步骤名","body":"说明..."}, ...],
            "image_keywords":"关键词"}

        5. two_column — 双栏布局
           {"layout":"two_column","title":"页面标题",
            "left_title":"左栏标题","left_items":["项1","项2"...],
            "right_title":"右栏标题","right_items":["项1","项2"...],
            "image_keywords":"关键词"}

        6. bullets — 要点列表
           {"layout":"bullets","title":"页面标题",
            "bullets":["要点1","要点2"...],
            "image_keywords":"关键词"}
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    # ── Theme colors ──
    if theme == "light":
        BG = RGBColor(0xF5, 0xF5, 0xF8)
        CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
        TITLE = RGBColor(0x1A, 0x1A, 0x2E)
        BODY = RGBColor(0x33, 0x33, 0x55)
        MUTED = RGBColor(0x88, 0x88, 0xAA)
        ACCENT1 = RGBColor(0x00, 0x74, 0xD9)
        ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
        LINE = RGBColor(0xDD, 0xDD, 0xEE)
    else:
        BG = RGBColor(0x1E, 0x1E, 0x3A)
        CARD_BG = RGBColor(0x25, 0x25, 0x45)
        TITLE = RGBColor(0xFF, 0xFF, 0xFF)
        BODY = RGBColor(0xDD, 0xDD, 0xEE)
        MUTED = RGBColor(0x99, 0x99, 0xBB)
        ACCENT1 = RGBColor(0x00, 0xA8, 0xFF)
        ACCENT2 = RGBColor(0x8B, 0x5C, 0xF6)
        LINE = RGBColor(0x33, 0x33, 0x55)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    raw = slides_json.strip()
    # Strip markdown code block fences if present
    if raw.startswith("```"):
        # Remove opening fence (```json, ```, etc.)
        first_newline = raw.find("\n")
        if first_newline != -1:
            raw = raw[first_newline + 1:]
        # Remove closing fence
        if raw.endswith("```"):
            raw = raw[:-3].strip()
        elif raw.endswith("```\n"):
            raw = raw[:-4].strip()
        elif "\n```" in raw:
            raw = raw[: raw.rfind("\n```")].strip()

    slides_data = json.loads(raw)

    # Support both formats:
    #   - array: [slide, slide, ...]
    #   - object: {"theme":"dark","slides":[slide, slide, ...]}
    if isinstance(slides_data, dict):
        if "theme" in slides_data and theme == "dark":
            theme = slides_data.get("theme", theme)
        slides_data = slides_data.get("slides", slides_data.get("slides_json", [slides_data]))

    if not isinstance(slides_data, list):
        return f"错误: slides_json 解析后应为数组，实际为 {type(slides_data).__name__}"

    def _add_bg(slide, color=None):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color or BG

    def _add_img(slide, kw, left, top, width, height=None, idx=0):
        if height is None:
            height = width * 0.75
        path = _download_image(kw, get_outputs_dir(), idx)
        if path and path.exists():
            try:
                slide.shapes.add_picture(str(path), left, top, width, height)
                return True
            except Exception:
                return False
        return False

    def _add_box(slide, l, t, w, h, text, size=20, color=None, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(l, t, w, h)
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color or BODY
        p.font.bold = bold
        p.alignment = align
        return box

    def _add_card(slide, l, t, w, h):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        return shape

    def _top_bar(slide, color):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    for si, slide_def in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _add_bg(slide)
        layout = slide_def.get("layout", "bullets")
        title = slide_def.get("title", "")
        kw = slide_def.get("image_keywords", "technology")
        accent = ACCENT1 if si % 2 == 0 else ACCENT2

        if layout == "cover":
            # Full image background on right half
            _add_img(slide, kw, Inches(7.0), Inches(0.5), Inches(6.0), Inches(6.5), si)
            _top_bar(slide, accent)
            _add_box(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(1.2),
                     title, 48, TITLE, True)
            # Decorative line
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(3.0), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = accent
            line.line.fill.background()
            # Subtitle
            sub = slide_def.get("subtitle", "")
            if sub:
                _add_box(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(0.6),
                         sub, 28, MUTED)
            # Items
            items = slide_def.get("items", [])
            for i, item in enumerate(items):
                _add_box(slide, Inches(0.8), Inches(4.2) + i * Inches(0.45),
                         Inches(5.5), Inches(0.4), item, 18, BODY)

        elif layout == "cards":
            _top_bar(slide, accent)
            _add_box(slide, Inches(0.5), Inches(0.3), Inches(1.0), Inches(0.5),
                     f"{si+1:02d}", 14, MUTED)
            _add_box(slide, Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.8),
                     title, 36, TITLE, True)
            _add_img(slide, kw, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5), si)

            cards = slide_def.get("cards", [])
            card_h = Inches(1.4)
            card_w = Inches(6.3)
            gap = Inches(0.2)
            start_y = Inches(1.7)

            for i, card in enumerate(cards):
                y = start_y + i * (card_h + gap)
                if i >= 6:
                    break
                _add_card(slide, Inches(0.5), y, card_w, card_h)
                _add_box(slide, Inches(0.8), y + Inches(0.08), card_w - Inches(0.6), Inches(0.35),
                         f"▸ {card.get('title', '')}", 16, accent, True)
                _add_box(slide, Inches(0.8), y + Inches(0.4), card_w - Inches(0.6), Inches(0.9),
                         card.get("body", ""), 12, BODY)

        elif layout == "grid":
            _top_bar(slide, accent)
            _add_box(slide, Inches(0.5), Inches(0.3), Inches(1.0), Inches(0.5),
                     f"{si+1:02d}", 14, MUTED)
            _add_box(slide, Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.8),
                     title, 36, TITLE, True)
            _add_img(slide, kw, Inches(9.5), Inches(6.0), Inches(3.3), Inches(1.3), si)

            cards = slide_def.get("cards", [])
            cw = Inches(5.8)
            ch = Inches(2.4)
            gap = Inches(0.3)
            sx = Inches(0.5)
            sy = Inches(1.7)

            for i, card in enumerate(cards):
                if i >= 4:
                    break
                col = i % 2
                row = i // 2
                x = sx + col * (cw + gap)
                y = sy + row * (ch + gap)
                _add_card(slide, x, y, cw, ch)
                _add_box(slide, x + Inches(0.3), y + Inches(0.15), cw - Inches(0.6), Inches(0.35),
                         f"✦ {card.get('title', '')}", 18, accent, True)
                _add_box(slide, x + Inches(0.3), y + Inches(0.6), cw - Inches(0.6), Inches(1.6),
                         card.get("body", ""), 12, BODY)

        elif layout == "workflow":
            _top_bar(slide, accent)
            _add_box(slide, Inches(0.5), Inches(0.3), Inches(1.0), Inches(0.5),
                     f"{si+1:02d}", 14, MUTED)
            _add_box(slide, Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.8),
                     title, 36, TITLE, True)

            steps = slide_def.get("steps", [])
            sw = Inches(2.8)
            sh = Inches(3.0)
            sx = Inches(0.5)
            sy = Inches(1.8)
            gap = Inches(0.35)

            for i, step in enumerate(steps):
                x = sx + i * (sw + gap)
                if i >= 5:
                    break
                # Step number circle
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), sy, Inches(0.45), Inches(0.45))
                circle.fill.solid()
                circle.fill.fore_color.rgb = accent
                circle.line.fill.background()
                tf = circle.text_frame
                tf.paragraphs[0].text = str(i + 1)
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = TITLE
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                _add_card(slide, x, sy + Inches(0.65), sw, sh - Inches(0.65))
                _add_box(slide, x + Inches(0.2), sy + Inches(0.8), sw - Inches(0.4), Inches(0.35),
                         step.get("title", ""), 18, accent, True)
                _add_box(slide, x + Inches(0.2), sy + Inches(1.25), sw - Inches(0.4), Inches(1.4),
                         step.get("body", ""), 12, BODY)

                # Arrow
                if i < len(steps) - 1:
                    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                                   x + sw + Inches(0.02), sy + Inches(1.3),
                                                   Inches(0.3), Inches(0.25))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = accent
                    arrow.line.fill.background()

            _add_img(slide, kw, Inches(0.5), Inches(5.3), Inches(12.3), Inches(2.0), si)

        elif layout == "two_column":
            _top_bar(slide, accent)
            _add_box(slide, Inches(0.5), Inches(0.3), Inches(1.0), Inches(0.5),
                     f"{si+1:02d}", 14, MUTED)
            _add_box(slide, Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.8),
                     title, 36, TITLE, True)

            l_title = slide_def.get("left_title", "")
            r_title = slide_def.get("right_title", "")
            l_items = slide_def.get("left_items", [])
            r_items = slide_def.get("right_items", [])

            _add_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                     l_title, 22, accent, True)
            for i, item in enumerate(l_items):
                _add_box(slide, Inches(0.5), Inches(2.2) + i * Inches(0.55),
                         Inches(6.0), Inches(0.5), f"▸ {item}", 14, BODY)

            _add_box(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
                     r_title, 22, accent, True)
            for i, item in enumerate(r_items):
                if isinstance(item, str):
                    _add_box(slide, Inches(6.8), Inches(2.2) + i * Inches(0.55),
                             Inches(6.0), Inches(0.5), f"✦ {item}", 14, BODY)
                elif isinstance(item, dict):
                    _add_box(slide, Inches(6.8), Inches(2.2) + i * Inches(0.55),
                             Inches(2.0), Inches(0.5), f"✦ {item.get('title','')}", 14, ACCENT1, True)
                    _add_box(slide, Inches(9.0), Inches(2.2) + i * Inches(0.55),
                             Inches(4.0), Inches(0.5), item.get("body", ""), 13, BODY)

            _add_img(slide, kw, Inches(0.5), Inches(5.2), Inches(5.0), Inches(2.0), si)

        else:  # bullets
            _top_bar(slide, accent)
            _add_box(slide, Inches(0.5), Inches(0.3), Inches(1.0), Inches(0.5),
                     f"{si+1:02d}", 14, MUTED)
            _add_box(slide, Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.8),
                     title, 36, TITLE, True)
            _add_img(slide, kw, Inches(0.5), Inches(6.0), Inches(4.0), Inches(1.2), si)

            bullets = slide_def.get("bullets", [])
            for i, b in enumerate(bullets):
                _add_box(slide, Inches(0.5), Inches(1.7) + i * Inches(0.6),
                         Inches(12.0), Inches(0.5), f"• {b}", 16, BODY)

    # ── Save ──
    outputs_dir = get_outputs_dir()
    safe_filename = os.path.basename(filename)
    if not safe_filename.endswith(".pptx"):
        safe_filename += ".pptx"
    out_path = outputs_dir / safe_filename
    prs.save(str(out_path))

    # Clean old cache files (older than 1 hour)
    cache_dir = outputs_dir / ".img_cache"
    if cache_dir.exists():
        now = __import__("time").time()
        for f in cache_dir.iterdir():
            if now - f.stat().st_mtime > 3600:
                try:
                    f.unlink()
                except Exception:
                    pass

    return f"PPT 已成功保存到: {out_path}（{os.path.getsize(out_path) / 1024:.0f} KB，{len(slides_data)} 页）"


# ── Keep backward-compatible simple tools ──

def write_temp_json(content: str, tag: str = "ppt") -> str:
    """将大段 JSON 内容写入临时文件，供后续工具读取。
    用于解决 LLM 无法将超长 JSON 作为函数参数传递的问题。

    Args:
        content (str): JSON 字符串内容（PPT slides JSON、Excel 数据等）。
        tag (str): 标签，用于区分不同格式的临时文件，如 'ppt'、'excel'。默认 'ppt'。

    Returns:
        临时文件的完整路径，后续可传给 save_pptx 的 json_file 参数。
    """
    outputs_dir = get_outputs_dir()
    temp_dir = outputs_dir / ".temp"
    temp_dir.mkdir(parents=True, exist_ok=True)
    import time as _time
    ts = int(_time.time())
    file_path = temp_dir / f"{tag}_{ts}.json"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)
    return str(file_path)


def save_pptx(filename: str, slides_json: str = "", json_file: str = "") -> str:
    """保存内容为 PowerPoint (.pptx) 文档，自动配图和精美排版。
    Args:
        filename (str): 文件名称，例如 '产品介绍.pptx'。
        slides_json (str): 幻灯片内容的 JSON 字符串（短内容时使用）。支持两种格式：
            - 数组格式: [{"layout":"cover","title":"标题"...}, ...]
            - 对象格式: {"theme":"dark","slides":[{"layout":"cover",...}, ...]}
        json_file (str): 临时 JSON 文件路径（长内容时使用，由 write_temp_json 生成）。
            如果提供了 json_file，则忽略 slides_json，从文件读取。
    """
    if json_file:
        fp = pathlib.Path(json_file)
        if fp.exists():
            slides_json = fp.read_text(encoding="utf-8")
            # Clean up temp file
            try:
                fp.unlink()
            except Exception:
                pass
    return create_pptx(filename, "dark", slides_json)


def save_document(filename: str, content: str) -> str:
    """保存纯文本或 Markdown 内容到本地文件中。"""
    outputs_dir = get_outputs_dir()
    safe_filename = os.path.basename(filename)
    file_path = outputs_dir / safe_filename
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)
    return f"文本/Markdown 文档已成功保存到: {file_path}"


def save_docx(filename: str, content: str) -> str:
    """保存内容为 Word (.docx) 文档（自动去除 markdown 符号）。"""
    try:
        from docx import Document
    except ImportError:
        return "错误: 缺少 python-docx 库，请安装。"
    outputs_dir = get_outputs_dir()
    safe_filename = os.path.basename(filename)
    if not safe_filename.endswith(".docx"):
        safe_filename += ".docx"
    file_path = outputs_dir / safe_filename

    import re
    # Strip common markdown formatting
    text = content
    # Remove markdown headings: ## text or ### text
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    # Remove bold/italic: **text** or *text*
    text = re.sub(r'\*{1,3}([^*]+)\*{1,3}', r'\1', text)
    # Remove markdown horizontal rules: --- or ***
    text = re.sub(r'^[\*\-]{3,}\s*$', '', text, flags=re.MULTILINE)
    # Remove markdown list markers: * text or - text
    text = re.sub(r'^[\*\-]\s+', '', text, flags=re.MULTILINE)
    # Remove inline code: `text`
    text = re.sub(r'`([^`]+)`', r'\1', text)

    doc = Document()
    for paragraph in text.split("\n"):
        stripped = paragraph.strip()
        if stripped:
            doc.add_paragraph(stripped)
    doc.save(file_path)
    return f"Word 文档已成功保存到: {file_path}"


def save_excel(filename: str, data_json: str) -> str:
    """保存内容为 Excel (.xlsx) 文档。"""
    try:
        import openpyxl
    except ImportError:
        return "错误: 缺少 openpyxl 库，请安装。"
    outputs_dir = get_outputs_dir()
    safe_filename = os.path.basename(filename)
    if not safe_filename.endswith(".xlsx"):
        safe_filename += ".xlsx"
    file_path = outputs_dir / safe_filename
    try:
        data = json.loads(data_json)
    except json.JSONDecodeError:
        return "错误: data_json 无法解析为有效的 JSON。"
    wb = openpyxl.Workbook()
    ws = wb.active
    for row in data:
        if isinstance(row, list):
            ws.append(row)
    wb.save(file_path)
    return f"Excel 文档已成功保存到: {file_path}"
